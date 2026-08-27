"""
pptx_remediate.py — Add remediation suggestions to a .pptx via speaker notes.

CAPABILITY NOTE (verified against python-pptx 1.0.2):
  - python-pptx does NOT expose a comment API (Slide.add_comment / Slide.comments
    are absent). True PowerPoint comments (comments.xml + commentAuthors.xml parts)
    are not writable without hand-rolling the OPC parts.
  - PowerPoint has NO "track changes" concept at all (revision tracking is a Word
    feature). So for .pptx the reliable, native mechanism is SPEAKER NOTES.

Strategy used here: append remediation suggestions to each slide's speaker notes,
and additionally add a final "Suggested Changes Summary" slide listing all edits so
the changes are visible on the slides themselves (not only in notes view).

Usage:
    from pptx_remediate import add_note, add_summary_slide
    prs = Presentation("input.pptx")
    add_note(prs.slides[2], "Add alt text to the chart on this slide.")
    add_summary_slide(prs, ["Slide 3: add alt text to chart", "Slide 5: fix contrast"])
    prs.save("output_remediated.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def add_note(slide, text):
    """Append a remediation suggestion to a slide's speaker notes."""
    notes = slide.notes_slide
    pf = notes.notes_text_frame
    if pf.text:
        pf.text += "\n"
    pf.text += "REMEDIATION: " + text


def add_summary_slide(prs, suggestions, title="Suggested Changes Summary"):
    """Add a final slide listing all remediation suggestions inline."""
    blank = prs.slide_layouts[6]  # usually the 'blank' layout
    slide = prs.slides.add_slide(blank)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    btf = body.text_frame
    btf.word_wrap = True
    for i, s in enumerate(suggestions, 1):
        para = btf.paragraphs[0] if i == 1 else btf.add_paragraph()
        para.text = f"{i}. {s}"
        para.font.size = Pt(16)
    return slide


if __name__ == "__main__":
    prs = Presentation()
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[5])
    add_note(prs.slides[0], "Add alt text to any images on this slide.")
    add_summary_slide(
        prs,
        ["Slide 1: add alt text to images", "Slide 2: verify reading order"],
    )
    prs.save("/tmp/remediated_example.pptx")
    print("Wrote /tmp/remediated_example.pptx with speaker notes + summary slide.")